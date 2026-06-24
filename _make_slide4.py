"""슬라이드 4(GPTQ 선정) 단독 재제작 — 출처 명시 + '성숙한'→'검증된' 표현 수정."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
GRAY = RGBColor(0x5D, 0x6D, 0x7E)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x88, 0x49)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xEA, 0xF2, 0xF8)
CITE = RGBColor(0x7F, 0x8C, 0x8D)
FONT = 'Malgun Gothic'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
s = prs.slides.add_slide(prs.slide_layouts[6])


def rect(x, y, w, h, c):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = c
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def txt(x, y, w, h, runs, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t, sz, c, b) in line:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = c
            r.font.bold = b; r.font.name = FONT
    return tb


rect(0, 0, SW, Inches(1.05), NAVY)
txt(Inches(0.5), Inches(0.22), Inches(11), Inches(0.7),
    [[('양자화 기법 선정 — 왜 GPTQ인가', 26, WHITE, True)]])
txt(Inches(12.3), Inches(0.32), Inches(0.8), Inches(0.5),
    [[('4', 16, RGBColor(0x9B, 0xB5, 0xD4), True)]], PP_ALIGN.RIGHT)

txt(Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.5),
    [[('두 기둥으로 결정: ', 15, NAVY, True), ('① 우리가 깎는 부분에 알고리즘이 맞는가  ', 15, BLUE, True),
      ('② 그 모델을 처리하는 도구가 있는가', 15, BLUE, True)]])

rect(Inches(0.6), Inches(1.75), Inches(12.13), Inches(0.05), BLUE)
txt(Inches(0.6), Inches(1.85), Inches(12), Inches(0.45),
    [[('근거 1 — 알고리즘 적합성 (출처 있음)', 17, BLUE, True)]])
txt(Inches(0.85), Inches(2.42), Inches(11.9), Inches(2.3), [
    [('· GPTQ = dense transformer의 Linear를 레이어별 Hessian 기반으로 양자화하는 기법', 15, NAVY, False)],
    [('     ', 11, GRAY, False), ('[GPTQ 원논문, arXiv:2210.17323]', 12, CITE, False)],
    [('· 우리가 4bit로 깎는 mllama 텍스트 디코더 = 정통 Llama dense 구조 (MoE·희소구조 없음)', 15, NAVY, False)],
    [('     ', 11, GRAY, False), ('[modeling_mllama.py 코드 직접 확인 · config model_type=mllama_text_model]', 12, CITE, False)],
    [('· Llama 계열 INT4(group_size=128): fp16 대비 평균 약 97% 정확도 유지', 15, NAVY, False)],
    [('     ', 11, GRAY, False), ('[Llama-3.1-8B INT4 예: MMLU 70.25 / 평균 97% recovery — LLMC, arXiv:2405.06001]', 12, CITE, False)],
])

rect(Inches(0.6), Inches(4.75), Inches(12.13), Inches(0.05), BLUE)
txt(Inches(0.6), Inches(4.85), Inches(12), Inches(0.45),
    [[('근거 2 — 도구 적용성', 17, BLUE, True)]])
txt(Inches(0.85), Inches(5.42), Inches(11.9), Inches(1.1), [
    [('· mllama(비전+텍스트 멀티모달) 구조를 처리하는 검증된 도구 = gptqmodel', 15, NAVY, False)],
    [('· AWQ·SpinQuant는 mllama 처리 로직이 없어 직접 구현 필요 → 도구 개발 수준의 비용', 15, NAVY, False)],
    [('     ', 11, GRAY, False),
     ('[AutoAWQ 유지보수 중단 · SpinQuant/llm-awq mllama 미지원 — 단, 알고리즘상 불가능이 아니라 구현 부재]', 12, CITE, False)],
])

rect(Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.62), PALE)
txt(Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5),
    [[('결론: ', 14, RED, True),
      ('Llama 구조라 GPTQ가 잘 맞고(①, 수치 입증) + 그 mllama를 실제로 처리하는 도구가 gptqmodel(②) → 교집합이 GPTQ', 14, NAVY, True)]])

s.notes_slide.notes_text_frame.text = (
    'GPTQ를 고른 근거는 두 기둥입니다. 순서가 중요한데, 알고리즘 적합성(①)이 먼저고 도구(②)가 그다음입니다.\n\n'
    '근거 1, 알고리즘 적합성. GPTQ는 dense transformer의 Linear 레이어를 캘리브 통계(Hessian)로 한 층씩 양자화하는 기법입니다[arXiv:2210.17323]. '
    '그런데 우리가 4bit로 깎는 부분이 바로 mllama의 텍스트 디코더이고, 이건 MoE나 희소구조 없는 정통 Llama dense 구조입니다(코드로 직접 확인). '
    '즉 GPTQ가 가장 잘 다루는 대상이 정확히 우리 양자화 타깃입니다. 실제로 Llama 계열은 INT4 group_size 128에서 fp16 대비 평균 약 97% 정확도를 유지한다고 여러 벤치마크가 보고합니다[LLMC, arXiv:2405.06001]. 즉 "Llama라서 잘 맞는다"는 통념이 아니라 수치로 뒷받침됩니다.\n\n'
    '근거 2, 도구 적용성. 알고리즘이 맞아도 그걸 우리 모델에 적용해줄 도구가 필요합니다. mllama는 비전+텍스트가 섞인 특수 구조라 대부분의 양자화 도구엔 처리 로직이 없습니다. '
    'gptqmodel만 mllama를 지원했고(그마저 살짝 낡아 한 줄 수정해 썼습니다), AWQ·SpinQuant는 직접 구현해야 했습니다. '
    '여기서 정직하게 짚을 점 — 이건 "다른 도구로는 불가능"이 아니라 "구현이 안 돼 있어 직접 만들어야 한다"는 비용 문제입니다. 우리 목표는 도구 개발이 아니라 모델 배포이므로, 바퀴를 재발명하지 않는 게 합리적이었습니다.\n\n'
    '정리하면, Llama 구조라 알고리즘이 잘 맞고(수치 입증) + 그 mllama를 실제로 처리하는 검증된 도구가 gptqmodel이라는 두 근거의 교집합이 GPTQ였습니다.'
)

prs.save('slide4_GPTQ_selection.pptx')
print('saved')
