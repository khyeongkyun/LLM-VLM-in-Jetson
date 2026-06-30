"""오늘(2026-06-24) 작업 내용 PPT 생성.

한국어 캘리브레이션 GPTQ 재양자화 + PPL/K-DTCBench 평가 결과 발표 자료.
results/*.png 차트를 슬라이드에 삽입.

실행:
  python src/make_ppt.py
출력:
  results/report_2026-06-24.pptx
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"

# ── 색상 (차트와 통일) ──
NAVY   = RGBColor(0x2C, 0x3E, 0x50)
BLUE   = RGBColor(0x4C, 0x72, 0xB0)
GREEN  = RGBColor(0x55, 0xA8, 0x68)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
GRAY   = RGBColor(0x60, 0x60, 0x60)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xC4, 0x4E, 0x52)

FONT = "Malgun Gothic"

# 16:9
SW, SH = Inches(13.333), Inches(7.5)


def _set_font(run, size, bold=False, color=NAVY, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) 또는 (text, size, bold, color, level)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[:4]
        level = spec[4] if len(spec) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        _set_font(run, size, bold, color)
    return tb


def add_band(slide, top, height, color):
    """가로 색 띠."""
    shp = slide.shapes.add_shape(1, 0, top, SW, height)  # 1 = rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_proportion_bar(slide, left, top, width, height, segments):
    """가로 비율 막대. segments: list of (label, frac, color)."""
    x = left
    for label, frac, color in segments:
        seg_w = Emu(int(width * frac))
        shp = slide.shapes.add_shape(1, x, top, seg_w, height)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.color.rgb = WHITE; shp.line.width = Pt(1.5)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        _set_font(r, 14, True, WHITE)
        x = Emu(int(x) + int(seg_w))


def add_bullet_slide(prs, title, bullets, subtitle=None):
    """제목 + 불릿 슬라이드. bullets: list of (text, level, color?)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 상단 타이틀 띠
    add_band(slide, 0, Inches(1.15), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.85),
                [(title, 30, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5),
                    [(subtitle, 15, False, GRAY)])
    # 불릿
    top = Inches(2.0) if subtitle else Inches(1.55)
    tb = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), SH - top - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        text, level = b[0], b[1]
        color = b[2] if len(b) > 2 else NAVY
        size = 22 if level == 0 else 18
        bold = level == 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        prefix = "■  " if level == 0 else "–  "
        run = p.add_run()
        run.text = prefix + text
        _set_font(run, size, bold, color)
    return slide


def add_image_slide(prs, title, img_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(1.0), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.75),
                [(title, 26, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

    img_path = Path(img_path)
    if img_path.exists():
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size
        # 가용 영역
        avail_w = SW - Inches(1.0)
        avail_top = Inches(1.2)
        avail_h = SH - avail_top - (Inches(0.7) if caption else Inches(0.3))
        ratio = min(avail_w / iw, avail_h / ih)
        w = Emu(int(iw * ratio))
        h = Emu(int(ih * ratio))
        left = Emu(int((SW - w) / 2))
        slide.shapes.add_picture(str(img_path), left, avail_top, width=w, height=h)
        if caption:
            add_textbox(slide, Inches(0.6), SH - Inches(0.6), Inches(12.1), Inches(0.5),
                        [(caption, 14, False, GRAY)], align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    [(f"[이미지 없음: {img_path.name}]", 18, False, RED)])
    return slide


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, SH, NAVY)
    # 그린 액센트 바
    bar = slide.shapes.add_shape(1, Inches(0.9), Inches(2.55), Inches(0.18), Inches(2.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
    bar.shadow.inherit = False

    add_textbox(slide, Inches(1.3), Inches(2.5), Inches(11), Inches(2.2), [
        ("Llama 3.2-11B Vision  GPTQ 4-bit 양자화", 36, True, WHITE),
        ("한국어 캘리브레이션 개선 & 평가", 36, True, GREEN),
    ])
    add_textbox(slide, Inches(1.32), Inches(4.7), Inches(11), Inches(1.2), [
        ("한국어 PPL 손상 +19.6% → +6.6% 달성 · K-DTCBench 한국어 VLM 평가 구축", 18, False, RGBColor(0xC8,0xD2,0xDC)),
        ("2026-06-24", 16, False, RGBColor(0x90,0xA0,0xB0)),
    ])
    return slide


def add_table_slide(prs, ppl, bench):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(1.0), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.75),
                [("종합 결과 요약", 26, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

    headers = ["모델", "EN PPL", "KO PPL", "EN Δ", "KO Δ", "K-DTCBench"]
    fp16_en, fp16_ko = ppl["fp16"]["en"]["ppl"], ppl["fp16"]["ko"]["ppl"]

    def row(key, label, bench_key):
        en, ko = ppl[key]["en"]["ppl"], ppl[key]["ko"]["ppl"]
        d_en = "—" if key == "fp16" else f"+{(en/fp16_en-1)*100:.1f}%"
        d_ko = "—" if key == "fp16" else f"+{(ko/fp16_ko-1)*100:.1f}%"
        dtc = f'{bench[bench_key]["total_accuracy"]:.1%}' if bench_key and bench_key in bench else "—"
        return [label, f"{en:.2f}", f"{ko:.2f}", d_en, d_ko, dtc]

    rows = [
        row("fp16", "fp16 원본", "fp16"),
        row("gptq_en", "GPTQ 영어캘리브", None),
        row("gptq_ko", "GPTQ 한국어캘리브 ★", "kocalib"),
    ]

    nrows, ncols = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(
        nrows, ncols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.6))
    table = tbl_shape.table

    widths = [3.2, 1.7, 1.7, 1.7, 1.7, 1.7]
    for c, w in enumerate(widths):
        table.columns[c].width = Inches(w)

    # 헤더
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; _set_font(r, 15, True, WHITE)

    # 데이터
    for ri, rdata in enumerate(rows, start=1):
        highlight = ri == 3
        for c, val in enumerate(rdata):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9) if highlight else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            color = NAVY
            if c == 4 and val != "—":  # KO Δ 강조
                color = GREEN if highlight else RED
            _set_font(r, 14, highlight, color)

    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), [
        ("핵심 결론", 18, True, NAVY),
        ("한국어 혼합 캘리브레이션으로 한국어 PPL 손상을 +19.6% → +6.6%로 1/3 수준 감소", 16, False, GREEN, 1),
        ("영어 성능도 유지(+7.3% → +6.4%) — 한쪽을 희생하지 않음", 16, False, NAVY, 1),
        ("K-DTCBench 한국어 VLM 평가 파이프라인 신규 구축 (240문제, 재현 가능한 정확매칭 채점)", 16, False, NAVY, 1),
    ])
    return slide


def add_calib_slide(prs):
    """캘리브레이션 데이터셋 상세 설명."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(1.15), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.85),
                [("캘리브레이션 데이터셋 상세", 30, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.6), Inches(1.28), Inches(12.1), Inches(0.5),
                [("GPTQ는 소수의 대표 샘플로 레이어별 양자화 오차를 보정한다 — 그래서 '무엇을 보여주느냐'가 성능을 가른다", 14, False, GRAY)])

    # 본문 불릿
    bullets = [
        ("캘리브레이션이란?", 0),
        ("GPTQ는 가중치를 4-bit로 압축하며, 대표 텍스트 샘플을 흘려보내 레이어 출력 오차(Hessian)를 계산해 보정", 1),
        ("→ 캘리브 데이터가 닿는 분포는 정밀하게, 닿지 않는 분포는 거칠게 양자화됨", 1, RED),
        ("우리 구성: Wikipedia-KO 70% + Flickr30k-EN 30%", 0),
        ("Wikipedia-KO (wikimedia/wikipedia · 20231101.ko): 배포 도메인인 한국어 문어체 텍스트를 대표", 1),
        ("Flickr30k-EN (lmms-lab/flickr30k): 비전–언어 정렬 경로를 유지하기 위한 멀티모달 캡션", 1),
        ("왜 70/30 비율인가?", 0),
        ("핵심 배포 태스크가 한국어(회의록·문서)이므로 한국어 비중을 높임", 1),
        ("단, 비전 인코더와 연결된 캡션 경로도 일부 유지해야 멀티모달 성능이 보존됨 → 영어 30%", 1),
        ("왜 이 데이터셋인가?", 0),
        ("parquet 포맷이라 HF 스트리밍이 안정적 (이전 ZIP 데이터셋은 76,000+ 요청 무한루프)", 1),
        ("평가 벤치마크와 분리 → 과적합 방지 / PPL 평가셋과 같은 분포(in-distribution) → 손상 측정 신뢰도↑", 1),
    ]
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(12), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        text, level = b[0], b[1]
        color = b[2] if len(b) > 2 else NAVY
        size = 19 if level == 0 else 15
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(7)
        prefix = "■  " if level == 0 else "–  "
        r = p.add_run(); r.text = prefix + text
        _set_font(r, size, level == 0, color)

    # 하단 비율 막대
    add_textbox(slide, Inches(0.7), Inches(6.45), Inches(4), Inches(0.4),
                [("캘리브 데이터 구성 비율", 13, True, NAVY)])
    add_proportion_bar(slide, Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.45), [
        ("Wikipedia-KO  70%", 0.70, GREEN),
        ("Flickr30k-EN  30%", 0.30, ORANGE),
    ])
    return slide


def add_benchmark_slide(prs):
    """K-DTCBench 벤치마크 상세 + 후보 비교."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, Inches(1.15), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.85),
                [("평가 벤치마크 상세 — K-DTCBench", 30, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.6), Inches(1.28), Inches(12.1), Inches(0.5),
                [("PPL은 '언어 모델링' 손상을, 벤치마크는 '실제 태스크' 손상을 측정한다 — 두 지표가 상호 보완", 14, False, GRAY)])

    bullets = [
        ("K-DTCBench란?", 0),
        ("NCSOFT 공개, VARCO-VISION 논문(arXiv:2411.19103)에서 제안된 한국어 VLM 벤치마크", 1),
        ("한국어 네이티브 Document·Table·Chart 이미지 240문제 (카테고리당 80문제, 디지털 50%+손글씨 50%)", 1),
        ("4지선다 객관식(MCQA) — 모델이 A/B/C/D 중 하나를 고르고 정답과 정확매칭", 1),
        ("왜 K-DTCBench를 1순위로 선택했나?", 0),
        ("배포 유스케이스(회의록·문서 처리)와 직결 — 문서·표·차트 이해가 핵심 역량", 1, GREEN),
        ("정확매칭 채점 → LLM judge 불필요 → 빠르고 100% 재현 가능", 1, GREEN),
        ("한국어 네이티브(번역본 아님)라 실제 한국어 성능을 반영 / CC-BY-NC-4.0 연구 사용 가능", 1, GREEN),
    ]
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(12), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        text, level = b[0], b[1]
        color = b[2] if len(b) > 2 else NAVY
        size = 19 if level == 0 else 15
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(6)
        prefix = "■  " if level == 0 else "–  "
        r = p.add_run(); r.text = prefix + text
        _set_font(r, size, level == 0, color)

    # 후보 비교 표
    add_textbox(slide, Inches(0.7), Inches(4.95), Inches(8), Inches(0.4),
                [("후보 벤치마크 비교 및 선정 결과", 15, True, NAVY)])
    cand = [
        ["벤치마크", "형식", "채점 방식", "선정"],
        ["K-DTCBench", "4지선다 MCQA", "정확매칭 (judge 불필요)", "★ 1순위"],
        ["KOFFVQA", "자유생성 VQA", "gemma-2-9b-it judge", "2순위 (보완)"],
        ["K-MMBench / K-MMStar", "MCQA", "—", "제외 (영어 번역본)"],
        ["TextVQA", "VQA", "—", "제외 (영어 전용)"],
    ]
    nrows, ncols = len(cand), 4
    tbl_shape = slide.shapes.add_table(nrows, ncols, Inches(0.7), Inches(5.35),
                                       Inches(11.9), Inches(1.85))
    table = tbl_shape.table
    for c, w in enumerate([3.4, 2.7, 3.4, 2.4]):
        table.columns[c].width = Inches(w)
    for c, h in enumerate(cand[0]):
        cell = table.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; _set_font(r, 13, True, WHITE)
    for ri, rdata in enumerate(cand[1:], start=1):
        chosen = ri == 1
        for c, val in enumerate(rdata):
            cell = table.cell(ri, c); cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9) if chosen else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c != 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            color = GREEN if (chosen and c == 3) else (GRAY if val.startswith("제외") else NAVY)
            _set_font(r, 12, chosen, color)
    return slide


def main():
    ppl = {
        "fp16": json.loads((RESULTS / "ppl_fp16.json").read_text()),
        "gptq_en": json.loads((RESULTS / "ppl_gptq.json").read_text()),
        "gptq_ko": json.loads((RESULTS / "ppl_gptq_kocalib.json").read_text()),
    }
    bench = {}
    if (RESULTS / "kdtcbench_kocalib.json").exists():
        bench["kocalib"] = json.loads((RESULTS / "kdtcbench_kocalib.json").read_text())
    for fp16_name in ("kdtcbench_fp16.json", "kdtcbench_fp16_est.json"):  # full 240 우선
        if (RESULTS / fp16_name).exists():
            bench["fp16"] = json.loads((RESULTS / fp16_name).read_text())
            break

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # 1. 타이틀
    add_title_slide(prs)

    # 2. 오늘의 작업 요약
    add_bullet_slide(prs, "오늘의 작업 요약", [
        ("한국어 혼합 캘리브레이션으로 GPTQ 4-bit 재양자화 완료", 0),
        ("Wikipedia-KO 70% + Flickr30k-EN 30% 캘리브 데이터 구성", 1),
        ("PPL 재측정 — 한국어 손상도 대폭 개선 확인", 0),
        ("KO Δ +19.6% → +6.6%, EN Δ +7.3% → +6.4%", 1, GREEN),
        ("K-DTCBench 한국어 VLM 평가 파이프라인 신규 구축 및 실행", 0),
        ("한국어 Document/Table/Chart 240문제 4지선다 VQA", 1),
        ("재양자화 과정의 버그 3종 수정", 0),
        ("HF ZIP 스트리밍 무한루프 / 캘리브 샘플 OOM / processor 파일 누락", 1),
        ("결과 시각화 자료 생성 (PPL · K-DTCBench · 종합 대시보드)", 0),
    ], subtitle="2026-06-24 · Llama 3.2-11B Vision 경량화 (Jetson Orin Nano 타깃)")

    # 3. 문제 정의
    add_bullet_slide(prs, "문제 — 영어 캘리브의 한국어 손상", [
        ("기존 양자화는 flickr30k 영어 캡션만으로 캘리브레이션", 0),
        ("한국어 경로가 '거의 안 쓰임'으로 분류되어 거칠게 양자화됨", 1, RED),
        ("PPL 측정 결과 언어 간 손상 불균형 확인", 0),
        ("한국어 PPL +19.6% — 영어(+7.3%)의 약 2.7배", 1, RED),
        ("배포 도메인은 한국어 태스크(회의록) → 손상이 치명적", 1),
        ("원인: 캘리브 도메인(영어)과 배포 도메인(한국어) 불일치", 0),
    ], subtitle="GPTQ는 캘리브 데이터가 닿지 않는 경로를 거칠게 근사한다")

    # 4. 해결 방법
    add_bullet_slide(prs, "해결 — 한국어 혼합 캘리브 재양자화", [
        ("캘리브레이션 데이터를 배포 도메인에 맞게 재구성", 0),
        ("Wikipedia-KO 70% (텍스트) + Flickr30k-EN 30% (멀티모달)", 1, BLUE),
        ("평가 벤치마크와 분리하여 과적합 방지", 1),
        ("재양자화 수행 — 텍스트 디코더만 4-bit, 비전 인코더는 fp16 유지", 0),
        ("최종 모델 크기 약 11GB, 무결성 검증 PASS", 1, GREEN),
        ("재양자화 중 발견·수정한 버그", 0),
        ("HF ZIP 데이터셋 스트리밍 76,000+ HTTP 요청 무한루프 → parquet으로 교체", 1),
        ("긴 위키 문서로 인한 CUDA OOM(361GB 할당) → 텍스트 길이 절단", 1),
        ("GPTQ 저장 시 image processor 파일 누락 → fp16에서 복사", 1),
    ], subtitle="배포 도메인을 닮은 캘리브 데이터로 한국어 경로를 '쓰이게' 만든다")

    # 5. 캘리브레이션 데이터셋 상세
    add_calib_slide(prs)

    # 6. PPL 결과
    add_image_slide(prs, "결과 ① — PPL 비교", RESULTS / "report_ppl.png",
                    caption="한국어 PPL 손상 +19.6% → +6.6% (약 1/3 수준). 영어 PPL도 +7.3% → +6.4%로 소폭 개선.")

    # 7. 벤치마크 상세
    add_benchmark_slide(prs)

    # 8. K-DTCBench
    add_image_slide(prs, "결과 ② — K-DTCBench 한국어 VLM 평가", RESULTS / "report_kdtcbench.png",
                    caption="핵심: fp16 원본도 29.2%(≈랜덤 25%) → 낮은 점수는 양자화 손상이 아니라 모델 한계(Llama 3.2 Vision은 영어 전용). 비전 경로는 양쪽 다 fp16 → 양자화 영향 노이즈 수준. (240문제 동일 조건, fp16 vs 한국어캘리브)")

    # 7. 종합 대시보드
    add_image_slide(prs, "결과 ③ — 종합 대시보드", RESULTS / "report_summary.png")

    # 8. 요약 표 + 결론
    add_table_slide(prs, ppl, bench)

    # 8.5 압축 시도 (3-bit · depth pruning)
    add_image_slide(prs, "추가 실험 — 8GB 적재를 위한 압축 시도", RESULTS / "report_compression.png",
                    caption="3-bit·depth pruning 모두 8GB 미달: 크기 병목은 양자화 안 된 fp16 비전 경로(~4.5GB). 8GB 적재는 비전경로 양자화 + 프루닝 LoRA 힐링 + sub-4bit 조합(Phase-2)이라야.")

    # 9. 다음 단계
    add_bullet_slide(prs, "다음 단계", [
        ("fp16 기준선 K-DTCBench 240문제 측정 완료 — 전 모델 동일 조건 확보", 0),
        ("결론: 양자화 손상은 노이즈 수준, 낮은 점수는 모델(영어전용) 한계", 1, GREEN),
        ("8GB Jetson 적재 = Phase-2 (one-shot 3-bit·프루닝으론 불가 실증)", 0),
        ("비전/cross-attn 양자화(멀티모달 캘리브) + 프루닝 LoRA 힐링 + sub-4bit(AQLM)", 1, BLUE),
        ("대안: SmolVLM2-2.2B 등 소형 VLM (4-bit ~1.3GB로 8GB 여유)", 1),
        ("KOFFVQA 2순위 벤치마크 — 자유서술 한국어 VQA, LLM judge 채점", 0),
    ], subtitle="양자화 검증(완료) → 한국어 태스크 평가 → 엣지 배포(Phase-2)")

    out = RESULTS / "report_2026-06-24.pptx"
    prs.save(str(out))
    print(f"[done] {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
